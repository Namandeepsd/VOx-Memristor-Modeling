import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Use 16:9 layout
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Layouts
TITLE_SLIDE = 0
BULLET_SLIDE = 1
BLANK_SLIDE = 6

def add_title_slide(title, subtitle, info):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"{subtitle}\n\n{info}"
    
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = bullet_points[0]
    for bp in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 0
    return slide

def add_image_slide(title, image_paths, captions, positions):
    # positions is a list of tuples (left, top, width)
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE])
    
    # Add title manually
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    
    for img, cap, pos in zip(image_paths, captions, positions):
        if os.path.exists(img):
            slide.shapes.add_picture(img, Inches(pos[0]), Inches(pos[1]), width=Inches(pos[2]))
            
            # Add caption
            cxBox = slide.shapes.add_textbox(Inches(pos[0]), Inches(pos[1] + pos[3]), Inches(pos[2]), Inches(1.0))
            ctf = cxBox.text_frame
            cp = ctf.paragraphs[0]
            cp.text = cap
            cp.alignment = PP_ALIGN.CENTER
            cp.font.size = Pt(16)
        else:
            print(f"Missing image: {img}")

# 1. Title slide
add_title_slide(
    "VOx Memristor: First-Principles Electro-Thermal Compact Model",
    "Comprehensive Design, Derivation, and Validation",
    "Name: _________________ | Course: _________________ | Date: _________________"
)

# 2. Aims & Objectives
add_bullet_slide(
    "Aims & Objectives",
    [
        "Derive a 3-ODE system coupling thermodynamics, phase kinetics, and parasitics.",
        "Implement a modular Python simulator utilizing implicit numerical solvers (Radau).",
        "Reproduce key experimental observations (NDR, pinched hysteresis, thermal loops).",
        "Validate the model against data for both voltage-biased and current-biased modes."
    ]
)

# 3. Physical Mechanism
slide3 = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
slide3.shapes.title.text = "Physical Mechanism"
tf3 = slide3.shapes.placeholders[1].text_frame
tf3.text = "Voltage/Current → Joule Heating → Temperature Rise → Metal-Insulator Transition (MIT) → Metallic Growth → Resistance Change"
p3 = tf3.add_paragraph()
p3.text = "Key Transitions: T_IMT = 315 K, T_MIT = 298 K"

# 4. Governing Equations
slide4 = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
slide4.shapes.title.text = "Governing Equations"
tf4 = slide4.shapes.placeholders[1].text_frame
tf4.text = "Thermal Conservation (First Law of Thermodynamics):"
tf4.paragraphs[0].font.bold = True
p = tf4.add_paragraph()
p.text = "C_th * (dT/dt) = P_joule - G_th * (T - T_amb) - L * (dφ/dt)"
p.font.name = "Courier New"
p.level = 1

p = tf4.add_paragraph()
p.text = "Phase Kinetics (KJMA-motivated):"
p.font.bold = True
p = tf4.add_paragraph()
p.text = "dφ/dt = (φ_eq(T) - φ) / τ_φ"
p.font.name = "Courier New"
p.level = 1

p = tf4.add_paragraph()
p.text = "Resistance Model (Bruggeman Effective Medium Theory):"
p.font.bold = True
p = tf4.add_paragraph()
p.text = "R = f(R_insulating(T), R_metallic(T), φ)"
p.font.name = "Courier New"
p.level = 1

# 5. Model Evolution Overview
add_bullet_slide(
    "Model Evolution Overview",
    [
        "Ideal Model: Proof-of-concept for fundamental physical mechanism.",
        "Real Model: Adds non-ideal semiconductor physics (Arrhenius, TCR, contact resistance).",
        "Advanced Model: Adds Bruggeman geometric percolation and circuit parasitics."
    ]
)

# 6. Codebase Architecture
add_bullet_slide(
    "Codebase Architecture",
    [
        "vox_parameters.py: Central dataclass for physical and simulation parameters.",
        "vox_model.py: Core physics equations for the ideal 2-ODE model.",
        "vox_real_model.py: Non-ideal real model with Arrhenius conduction and metallic TCR.",
        "vox_advanced_model.py: 3-ODE system with Bruggeman EMT and parasitic capacitance.",
        "vox_solver.py: Radau implicit solver handling stiff system integration.",
        "vox_waveforms.py: Generators for voltage and current stimulation waveforms."
    ]
)

# 7. Ideal Model Results
add_image_slide(
    "Ideal Model Results",
    ["output/plots/dashboard.png"],
    ["The 4-panel dashboard demonstrates the fundamental phase transition and temperature cycling under voltage bias."],
    [(2.5, 1.5, 8.333, 5.5)] # left, top, width, height_offset_for_caption
)

# 8. Ideal Model: R-T and I-V
add_image_slide(
    "Ideal Model: R-T and I-V",
    ["output/plots/rt_curve.png", "output/plots/iv_curve.png"],
    ["R-T loop correctly captures the metal-insulator transition.", "I-V curve exhibits classic asymmetric pinched hysteresis."],
    [(0.5, 2.0, 6.0, 5.0), (6.833, 2.0, 6.0, 5.0)]
)

# 9. Real Model: Physics Corrections
add_bullet_slide(
    "Real Model: Physics Corrections",
    [
        "Arrhenius R_i: Exponentially dropping insulating resistance.",
        "TCR R_m: Metallic resistance increases with temperature.",
        "Contact Resistance: Imposes a minimum series resistance floor.",
        "Arrhenius τ_φ: Temperature-dependent phase relaxation time.",
        "Nonlinear G_th: Enhanced heat dissipation at high temperatures."
    ]
)

# 10. Real Model Results
add_image_slide(
    "Real Model Results",
    ["output/plots/real_fig2_RT.png", "output/plots/real_dashboard.png"],
    ["Curved, non-linear pre-switching behavior introduced by Arrhenius conduction.", "The Real Model properly bounds high-current and steady-state thermal behavior."],
    [(0.5, 2.0, 5.5, 4.5), (6.5, 2.0, 6.5, 4.5)]
)

# 11. Advanced Model: Why It Was Needed
add_bullet_slide(
    "Advanced Model: Why It Was Needed",
    [
        "Bruggeman Percolation: Replaces simplified mixing with geometric percolation thresholds.",
        "Circuit Parasitics: Introduces parasitic capacitance in a stiff 3-ODE system, producing sharp snap-backs."
    ]
)

# 12. Advanced Model Results: NDR
add_image_slide(
    "Advanced Model Results: NDR",
    ["output/plots/reproduced_fig3.png", "output/plots/adv_fig3_NDR.png"],
    ["Model closely reproduces the experimental S-shaped NDR curve.", "Phase kinetics lag triggers the sharp snap-back hysteresis under current bias."],
    [(0.5, 1.5, 5.5, 5.0), (6.5, 1.5, 5.5, 5.0)]
)

# 13. Advanced Model Results: Amplitude Sweeps
add_image_slide(
    "Advanced Model Results: Amplitude Sweeps",
    ["output/plots/adv_fig5_multi.png"],
    ["Amplitude-dependent progression effectively captures saturation at ~3-4V under varying current limits."],
    [(2.5, 1.5, 8.333, 5.5)]
)

# 14. Advanced Model Results: Cycle-to-Cycle Memory
add_image_slide(
    "Advanced Model Results: Cycle-to-Cycle Memory",
    ["output/plots/reproduced_fig6.png"],
    ["Cycle 1 traces a wider outer loop due to phase memory; cycles 2-10 cluster tightly inside."],
    [(3.5, 1.5, 6.333, 5.2)]
)

# 15. Advanced Dashboard
add_image_slide(
    "Advanced Dashboard",
    ["output/plots/adv_dashboard.png"],
    ["The 3-ODE Advanced Model incorporates percolation and parasitic capacitance to capture experimental measurement conditions."],
    [(2.5, 1.5, 8.333, 5.5)]
)

# 16. Quantitative Validation
slide16 = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE])
slide16.shapes.title.text = "Quantitative Validation"
# Add a table manually
rows = 4
cols = 3
left = Inches(1.0)
top = Inches(2.0)
width = Inches(11.333)
height = Inches(2.0)
table = slide16.shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(4.166)
table.columns[2].width = Inches(4.166)

headers = ["Metric", "Paper Observation", "Model Performance"]
data = [
    ["Voltage Saturation", "~3-4 V at 20 mA", "~3 V at 20 mA"],
    ["Hysteresis Width", "15-20 K span", "17 K span (315 K to 298 K)"],
    ["Switching Threshold", "NDR snap-back near ~0.5 mA", "Closely reproduces near ~0.5 mA"]
]
for col_idx, header in enumerate(headers):
    table.cell(0, col_idx).text = header
for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        table.cell(row_idx+1, col_idx).text = val

# 17. Known Limitations
add_bullet_slide(
    "Known Limitations",
    [
        "Resistance Span Trade-off: Model captures ~2.5 orders of magnitude vs. the paper's ~3 to prioritize accurate 2-4V switching thresholds.",
        "Geometric Simplifications: 0D lumped-element approach lacks spatial filament expansion, making mathematical transitions sharper than reality.",
        "Semiconducting Pre-Conduction: The ideal model lacks Arrhenius conduction, yielding a perfectly linear pre-switching curve."
    ]
)

# 18. Conclusion
add_bullet_slide(
    "Conclusion",
    [
        "Ideal Model: Validates the fundamental physical mechanism utilizing minimal assumptions.",
        "Real Model: Incorporates non-ideal semiconductor physics to achieve quantitative experimental agreement.",
        "Advanced Model: Leverages Bruggeman percolation and circuit parasitics to mimic non-linear measurement realities."
    ]
)

# 19. References
add_bullet_slide(
    "References",
    [
        'Rana, A.S., del Valle, J. et al., "Symmetric and tunable resistive switching in VO2 memristors," Scientific Reports 10, 3293 (2020), https://doi.org/10.1038/s41598-020-60373-z'
    ]
)

prs.save('VOx_Memristor_Presentation.pptx')
print("Presentation generated successfully.")
